import os
from pptx import Presentation

def verify_all_slides(pptx_path, template_path):
    prs = Presentation(pptx_path)
    tmpl = Presentation(template_path)
    
    # Get template logo details
    tmpl_logo = None
    for shape in tmpl.slides[0].shapes:
        if shape.name == "Freeform 8":
            tmpl_logo = {
                'name': shape.name,
                'left': shape.left / 914400,
                'top': shape.top / 914400,
                'width': shape.width / 914400,
                'height': shape.height / 914400
            }
            break
            
    print(f"Template Logo: {tmpl_logo}")
    print(f"Total slides: {len(prs.slides)}")
    
    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i+1} ---")
        
        # 1. Check logo
        logo_shape = None
        for shape in slide.shapes:
            if shape.name == "Freeform 8":
                logo_shape = shape
                break
        if logo_shape is not None:
            left = logo_shape.left / 914400
            top = logo_shape.top / 914400
            width = logo_shape.width / 914400
            height = logo_shape.height / 914400
            match = (
                abs(left - tmpl_logo['left']) < 0.001 and
                abs(top - tmpl_logo['top']) < 0.001 and
                abs(width - tmpl_logo['width']) < 0.001 and
                abs(height - tmpl_logo['height']) < 0.001
            )
            print(f"  Logo: Found '{logo_shape.name}' at ({left:.3f}, {top:.3f}, {width:.3f}, {height:.3f}) - Match template: {match}")
        else:
            print("  Logo: NOT FOUND")
            
        # 2. Check Nav Bar elements
        nav_elements = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                if txt in ["Part 1", "Part 2", "Part 3", "Problem Analysis", "Supply Chain Mapping", "Data & Insights", "Proposed Solutions", "Financial Impact", "Detail Roadmap"]:
                    left = shape.left / 914400
                    top = shape.top / 914400
                    width = shape.width / 914400
                    height = shape.height / 914400
                    nav_elements.append((shape.name, txt, f"({left:.3f}, {top:.3f}, {width:.3f}, {height:.3f})"))
        print(f"  Nav Elements found: {nav_elements}")
        
        # 3. Check for any other TextBoxes or elements in nav area (top part of the slide, e.g. top < 1.5 inches)
        nav_area_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                top = shape.top / 914400 if shape.top is not None else 0
                if top < 1.0 and shape.name != "Freeform 8":
                    txt = shape.text_frame.text.strip().replace('\n', ' ')
                    nav_area_text.append((shape.name, txt))
        print(f"  All texts in Nav Area (top < 1.0 inch): {nav_area_text}")

if __name__ == "__main__":
    verify_all_slides(
        "/Users/bunnypro/teamwork_projects/logage_slides/output.pptx",
        "/Users/bunnypro/Projects/LOGage2026/Red Modern Logistic Presentation.pptx"
    )
