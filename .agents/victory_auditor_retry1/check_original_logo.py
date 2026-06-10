import os
from pptx import Presentation

def inspect_template(pptx_path):
    prs = Presentation(pptx_path)
    print(f"Slide width: {prs.slide_width / 914400:.2f} inches")
    print(f"Slide height: {prs.slide_height / 914400:.2f} inches")
    print(f"Total slides: {len(prs.slides)}")
    print("-" * 50)
    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i+1} ---")
        for shape in slide.shapes:
            left = shape.left / 914400 if shape.left is not None else 0
            top = shape.top / 914400 if shape.top is not None else 0
            width = shape.width / 914400 if shape.width is not None else 0
            height = shape.height / 914400 if shape.height is not None else 0
            text = ""
            if shape.has_text_frame:
                text = shape.text_frame.text.strip().replace('\n', ' | ')
            print(f"Shape: '{shape.name}' | Type: {shape.shape_type} | Box: ({left:.3f}, {top:.3f}, {width:.3f}, {height:.3f})")
            if text:
                print(f"  Text: {text[:100]}")

if __name__ == "__main__":
    inspect_template("/Users/bunnypro/Projects/LOGage2026/Red Modern Logistic Presentation.pptx")
