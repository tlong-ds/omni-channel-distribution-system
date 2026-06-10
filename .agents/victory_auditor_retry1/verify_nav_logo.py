import os
from pptx import Presentation

def verify_output(pptx_path):
    prs = Presentation(pptx_path)
    print(f"Total slides: {len(prs.slides)}")
    print("-" * 50)
    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i+1} ---")
        logo_found = False
        nav_elements = []
        for shape in slide.shapes:
            left = shape.left / 914400 if shape.left is not None else 0
            top = shape.top / 914400 if shape.top is not None else 0
            width = shape.width / 914400 if shape.width is not None else 0
            height = shape.height / 914400 if shape.height is not None else 0
            
            # Check for Freeform 8 or logo position
            # Template logo position is (18.623, 0.484, 1.377, 1.167)
            is_logo = False
            if abs(left - 18.623) < 0.01 and abs(top - 0.484) < 0.01 and abs(width - 1.377) < 0.01 and abs(height - 1.167) < 0.01:
                is_logo = True
                logo_found = True
            
            text = ""
            if shape.has_text_frame:
                text = shape.text_frame.text.strip().replace('\n', ' ')
            
            if text in ["Part 1", "Part 2", "Part 3"]:
                nav_elements.append((shape.name, text, (left, top, width, height)))
            
            if is_logo:
                print(f"  LOGO shape: Name='{shape.name}', Type={shape.shape_type}, Box=({left:.3f}, {top:.3f}, {width:.3f}, {height:.3f})")
            elif text:
                print(f"  Shape: '{shape.name}' | Text: '{text[:80]}' | Box=({left:.3f}, {top:.3f}, {width:.3f}, {height:.3f})")
        
        print(f"  Logo found matching template coordinates: {logo_found}")
        print(f"  Nav elements found: {nav_elements}")

if __name__ == "__main__":
    verify_output("/Users/bunnypro/teamwork_projects/logage_slides/output.pptx")
