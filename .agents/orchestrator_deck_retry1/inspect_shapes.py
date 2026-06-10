import os
from pptx import Presentation
from pptx.util import Inches

def inspect_presentation(pptx_path):
    prs = Presentation(pptx_path)
    
    print(f"Slide width: {prs.slide_width / 914400:.2f} inches")
    print(f"Slide height: {prs.slide_height / 914400:.2f} inches")
    print("-" * 50)
    
    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i+1} ---")
        
        # Check shapes
        for shape in slide.shapes:
            shape_name = shape.name
            shape_type = shape.shape_type
            left = shape.left / 914400 if shape.left is not None else 0
            top = shape.top / 914400 if shape.top is not None else 0
            width = shape.width / 914400 if shape.width is not None else 0
            height = shape.height / 914400 if shape.height is not None else 0
            
            text = ""
            if shape.has_text_frame:
                text = shape.text_frame.text.strip().replace('\n', ' | ')
                
            print(f"Shape: '{shape_name}' | Type: {shape_type} | Box: ({left:.3f}, {top:.3f}, {width:.3f}, {height:.3f})")
            if text:
                print(f"  Text: {text[:100]}..." if len(text) > 100 else f"  Text: {text}")
                
            # If it's a picture, print some details
            if shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    img_name = shape.image.filename if hasattr(shape.image, 'filename') else "unknown"
                    print(f"  Picture filename: {img_name}")
                except Exception as e:
                    print(f"  Picture info error: {e}")

if __name__ == "__main__":
    pptx_path = "/Users/bunnypro/teamwork_projects/logage_slides/output.pptx"
    inspect_presentation(pptx_path)
