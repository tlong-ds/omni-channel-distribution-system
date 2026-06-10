import os
from pptx import Presentation
from pptx.dml.color import RGBColor

def get_color_str(color):
    if color.type == 1:  # RGB
        return str(color.rgb)
    elif color.type == 2:  # Theme
        return f"Theme Color: {color.theme_color}"
    else:
        return "Unknown/None"

def inspect_text_formatting(pptx_path):
    prs = Presentation(pptx_path)
    
    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i+1} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                print(f"Shape: '{shape.name}' ({shape.left/914400:.2f}, {shape.top/914400:.2f})")
                for p_idx, p in enumerate(shape.text_frame.paragraphs):
                    p_text = p.text.strip()
                    if not p_text:
                        continue
                    print(f"  Paragraph {p_idx}: '{p_text[:60]}...'")
                    for r_idx, r in enumerate(p.runs):
                        r_text = r.text.strip()
                        if not r_text:
                            continue
                        font = r.font
                        color_str = "None"
                        if font.color and font.color.type is not None:
                            color_str = get_color_str(font.color)
                        size_str = f"{font.size.pt:.1f}pt" if font.size else "None"
                        print(f"    Run {r_idx}: '{r_text}' | Font: {font.name} | Size: {size_str} | Color: {color_str} | Bold: {font.bold}")

if __name__ == "__main__":
    pptx_path = "/Users/bunnypro/teamwork_projects/logage_slides/output.pptx"
    inspect_text_formatting(pptx_path)
